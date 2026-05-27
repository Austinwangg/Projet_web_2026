"""
VoyageVista – PowerPoint Préliminaire v2
- Diagrammes MEA, SR et Architecture en PNG Mermaid
- Slides texte avec espacement corrigé
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pathlib

# ── Palette ─────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
CARD_BG    = RGBColor(0x16, 0x2A, 0x3E)
ACCENT1    = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2    = RGBColor(0x48, 0xCA, 0xB2)   # turquoise
ACCENT3    = RGBColor(0xFF, 0xA5, 0x00)   # orange
RED_ACC    = RGBColor(0xFF, 0x69, 0x69)
PURPLE_ACC = RGBColor(0x9B, 0x59, 0xB6)
GREEN_ACC  = RGBColor(0x2E, 0xCC, 0x71)
YELLOW_ACC = RGBColor(0xF3, 0x9C, 0x12)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xC8, 0xD6, 0xE5)
MID_GREY   = RGBColor(0x5A, 0x7A, 0x9A)
TEXT_DARK  = RGBColor(0x0D, 0x1B, 0x2A)

W = Inches(13.33)
H = Inches(7.5)
DIAG = pathlib.Path(r"C:\MAMP\htdocs\Projet_web_2026\diagrams")

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]
slide_n = [0]

# ── Helpers ──────────────────────────────────────────────────────────────────
def new_slide():
    slide_n[0] += 1
    return prs.slides.add_slide(BLANK)

def bg(slide, color=DARK_BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = color
    return s

def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_para(tf, text, size=12, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    p = tf.add_paragraph(); p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return p

def hline(slide, y, color=ACCENT1, thick=2):
    s = slide.shapes.add_shape(1, Inches(0.4), y, W - Inches(0.8), Pt(thick))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def section_tag(slide, label, color=ACCENT1):
    rect(slide, Inches(0.4), Inches(0.22), Inches(0.07), Inches(0.38), color)
    txt(slide, label, Inches(0.55), Inches(0.20), Inches(5), Inches(0.42),
        size=11, bold=True, color=color)

def slide_title(slide, title, subtitle=None):
    txt(slide, title, Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.72),
        size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.42),
            size=13, color=LIGHT_GREY, italic=True)

def pagenum(slide):
    txt(slide, str(slide_n[0]),
        W - Inches(0.6), H - Inches(0.4), Inches(0.45), Inches(0.3),
        size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)

def img(slide, png_name, l, t, w=None, h=None):
    path = str(DIAG / png_name)
    if w and h:
        slide.shapes.add_picture(path, l, t, w, h)
    elif w:
        slide.shapes.add_picture(path, l, t, width=w)
    elif h:
        slide.shapes.add_picture(path, l, t, height=h)
    else:
        slide.shapes.add_picture(path, l, t)

# ── SLIDE 1: PAGE DE GARDE ───────────────────────────────────────────────────
s = new_slide(); bg(s)
rect(s, 0, 0, Inches(4.6), H, CARD_BG)
rect(s, 0, 0, Inches(0.12), H, ACCENT1)

txt(s, "VOYAGEVISTA", Inches(0.55), Inches(1.4), Inches(3.9), Inches(0.75),
    size=30, bold=True, color=ACCENT1)
txt(s, "Plateforme de voyage en ligne",
    Inches(0.55), Inches(2.18), Inches(3.9), Inches(0.45),
    size=13, color=LIGHT_GREY, italic=True)
rect(s, Inches(0.55), Inches(2.7), Inches(3.5), Inches(0.04), ACCENT2)

# Zone à compléter
rect(s, Inches(0.55), Inches(2.9), Inches(3.7), Inches(2.5), RGBColor(0x0A,0x15,0x22))
for i, lbl in enumerate(["[ Noms des membres ]", "[ Groupe / Promotion ]",
                          "[ Date de rendu ]", "[ Enseignant ]"]):
    txt(s, lbl, Inches(0.72), Inches(3.05) + Inches(i * 0.55), Inches(3.4), Inches(0.42),
        size=12, color=MID_GREY, italic=True)

# Droite
txt(s, "PowerPoint Préliminaire", Inches(5.0), Inches(1.2), Inches(7.9), Inches(0.65),
    size=22, bold=True, color=WHITE)
txt(s, "Livrable de conception — Projet Web 2026",
    Inches(5.0), Inches(1.85), Inches(7.9), Inches(0.4),
    size=12, color=LIGHT_GREY, italic=True)

sections = [
    ("01", "Wireframes & Interfaces utilisateur"),
    ("02", "Storyboard & Parcours utilisateurs"),
    ("03", "Modèle Entité-Association"),
    ("04", "Schéma Relationnel"),
    ("05", "Architecture Système"),
    ("06", "Spécifications Fonctionnelles"),
]
for i, (num, label) in enumerate(sections):
    y = Inches(2.55) + i * Inches(0.65)
    rect(s, Inches(5.0), y, Inches(0.52), Inches(0.48), CARD_BG)
    txt(s, num, Inches(5.02), y + Inches(0.05), Inches(0.48), Inches(0.38),
        size=13, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    txt(s, label, Inches(5.62), y + Inches(0.08), Inches(7.1), Inches(0.38), size=13)

pagenum(s)

# ── SLIDE 2: SOMMAIRE ────────────────────────────────────────────────────────
s = new_slide(); bg(s)
section_tag(s, "SOMMAIRE")
slide_title(s, "Table des matières")
hline(s, Inches(1.5))

items = [
    ("01 — Wireframes & Interfaces",      "Slides 3–5",   ACCENT1),
    ("02 — Storyboard & Parcours",        "Slides 6–8",   ACCENT2),
    ("03 — Modèle Entité-Association",    "Slides 9–12",  ACCENT3),
    ("04 — Schéma Relationnel",           "Slides 13–15", RED_ACC),
    ("05 — Architecture Système",         "Slides 16–19", PURPLE_ACC),
    ("06 — Spécifications Fonctionnelles","Slides 20–25", GREEN_ACC),
]
for i, (label, pages, color) in enumerate(items):
    col = i % 2; row = i // 2
    x = Inches(0.4) + col * Inches(6.55)
    y = Inches(1.75) + row * Inches(1.6)
    rect(s, x, y, Inches(6.2), Inches(1.4), CARD_BG)
    rect(s, x, y, Inches(0.1), Inches(1.4), color)
    txt(s, label, x + Inches(0.25), y + Inches(0.2),
        Inches(5.7), Inches(0.6), size=15, bold=True)
    txt(s, pages, x + Inches(0.25), y + Inches(0.82),
        Inches(5.7), Inches(0.4), size=12, color=color, italic=True)

pagenum(s)

# ── SLIDES 3-5: WIREFRAMES (zones à compléter) ───────────────────────────────
wf_data = [
    ("Wireframes — Pages principales",
     "Accueil · Résultats · Détail destination",
     ["Page Accueil (Hero + Recherche + Destinations)", "Page Résultats (Filtres + Cartes)", "Page Détail Destination (4 onglets)"]),
    ("Wireframes — Espace utilisateur & Réservation",
     "Panier · Paiement · Compte · Itinéraire",
     ["Page Panier / Récapitulatif", "Page Paiement & Confirmation", "Espace Compte Utilisateur", "Créateur d'Itinéraire"]),
    ("Wireframes — Transport & Administration",
     "Recherche transport multimodale · Dashboard Admin",
     ["Page Recherche Transport (avion, train, bus, voiture)", "Dashboard Admin (gestion utilisateurs & données)"]),
]
for i, (title, sub, items) in enumerate(wf_data):
    s = new_slide(); bg(s)
    section_tag(s, "01 — WIREFRAMES")
    slide_title(s, title, sub)
    hline(s, Inches(1.55))

    rect(s, Inches(0.4), Inches(1.75), Inches(8.5), Inches(5.4), CARD_BG)
    rect(s, Inches(0.4), Inches(1.75), Inches(8.5), Inches(0.05), ACCENT1)
    txt(s, "[ Insérer wireframe ici ]",
        Inches(0.4), Inches(4.1), Inches(8.5), Inches(0.75),
        size=16, color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)

    rect(s, Inches(9.1), Inches(1.75), Inches(3.85), Inches(5.4), CARD_BG)
    rect(s, Inches(9.1), Inches(1.75), Inches(0.08), Inches(5.4), ACCENT1)
    txt(s, "Pages représentées",
        Inches(9.28), Inches(1.88), Inches(3.5), Inches(0.4),
        size=13, bold=True, color=ACCENT1)
    tb = s.shapes.add_textbox(Inches(9.28), Inches(2.38), Inches(3.55), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "   " + item
        r.font.size = Pt(12); r.font.color.rgb = LIGHT_GREY

    pagenum(s)

# ── SLIDES 6-8: STORYBOARD (zones à compléter) ──────────────────────────────
sb_data = [
    ("Storyboard — Parcours Réservation",
     "De la recherche à l'ajout au panier",
     ["Etape 1 : Accueil & Recherche", "Etape 2 : Consultation des résultats", "Etape 3 : Détail destination & sélection", "Etape 4 : Ajout au panier"]),
    ("Storyboard — Finalisation & Compte",
     "Paiement, confirmation, espace personnel",
     ["Etape 5 : Récapitulatif panier", "Etape 6 : Paiement & confirmation", "Etape 7 : Réception notification", "Etape 8 : Consultation espace compte"]),
    ("Storyboard — Création d'itinéraire",
     "Fonctionnalité de planification personnalisée",
     ["Etape 1 : Accès au créateur d'itinéraire", "Etape 2 : Ajout d'étapes & composants", "Etape 3 : Sauvegarde & consultation", "Etape 4 : Modification / Export"]),
]
for i, (title, sub, steps) in enumerate(sb_data):
    s = new_slide(); bg(s)
    section_tag(s, "02 — STORYBOARD")
    slide_title(s, title, sub)
    hline(s, Inches(1.55))

    rect(s, Inches(0.4), Inches(1.75), Inches(8.5), Inches(5.4), CARD_BG)
    rect(s, Inches(0.4), Inches(1.75), Inches(8.5), Inches(0.05), ACCENT2)
    txt(s, "[ Insérer storyboard / captures ici ]",
        Inches(0.4), Inches(4.1), Inches(8.5), Inches(0.75),
        size=16, color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)

    rect(s, Inches(9.1), Inches(1.75), Inches(3.85), Inches(5.4), CARD_BG)
    rect(s, Inches(9.1), Inches(1.75), Inches(0.08), Inches(5.4), ACCENT2)
    txt(s, "Etapes du parcours",
        Inches(9.28), Inches(1.88), Inches(3.5), Inches(0.4),
        size=13, bold=True, color=ACCENT2)
    tb = s.shapes.add_textbox(Inches(9.28), Inches(2.38), Inches(3.55), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    for step in steps:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "   " + step
        r.font.size = Pt(12); r.font.color.rgb = LIGHT_GREY

    pagenum(s)

# ── SLIDE 9: MEA — Vue d'ensemble ───────────────────────────────────────────
s = new_slide(); bg(s)
section_tag(s, "03 — MODELE ENTITE-ASSOCIATION")
slide_title(s, "MEA — Vue d'ensemble des entités",
            "6 modules fonctionnels · 14 entités · Relations identifiées")
hline(s, Inches(1.55))

modules = [
    ("Catalogue & Offres",   ["DESTINATION", "TRANSPORT", "HEBERGEMENT", "ACTIVITE"],       ACCENT1,    Inches(0.35),  Inches(1.72)),
    ("Reservations",         ["RESERVATION", "PAIEMENT", "RESERVATION_ACTIVITES"],          ACCENT2,    Inches(4.6),   Inches(1.72)),
    ("Itineraires",          ["ITINERAIRE", "ELEMENT_ITINERAIRE"],                           ACCENT3,    Inches(8.85),  Inches(1.72)),
    ("Utilisateurs",         ["UTILISATEUR", "SESSION", "NOTIFICATION"],                    RED_ACC,    Inches(0.35),  Inches(4.15)),
    ("Complements",          ["PANIER", "ELEMENT_PANIER", "FAVORI", "AVIS"],                PURPLE_ACC, Inches(4.6),   Inches(4.15)),
    ("Voyageur",             ["VOYAGEUR_SEJOUR"],                                            GREEN_ACC,  Inches(8.85),  Inches(4.15)),
]
for name, entities, color, x, y in modules:
    bh = Inches(0.55) + len(entities) * Inches(0.42)
    rect(s, x, y, Inches(3.95), bh, CARD_BG)
    rect(s, x, y, Inches(0.08), bh, color)
    txt(s, name, x + Inches(0.18), y + Inches(0.1),
        Inches(3.65), Inches(0.38), size=12, bold=True, color=color)
    for j, ent in enumerate(entities):
        txt(s, "  " + ent,
            x + Inches(0.18), y + Inches(0.52) + j * Inches(0.42),
            Inches(3.65), Inches(0.38), size=11, color=LIGHT_GREY)

pagenum(s)

# ── SLIDES 10-12: MEA diagrammes Mermaid ────────────────────────────────────
mea_slides = [
    ("MEA — Bloc Catalogue & Offres",
     "DESTINATION · TRANSPORT · HEBERGEMENT · ACTIVITE",
     "mea_catalogue.png",
     "DESTINATION (1) ---< TRANSPORT (N)  |  DESTINATION (1) ---< HEBERGEMENT (N)  |  DESTINATION (1) ---< ACTIVITE (N)\nToutes les entités offres référencent une destination via id_destination (FK)"),
    ("MEA — Reservations & Itineraires",
     "RESERVATION · PAIEMENT · ITINERAIRE · ELEMENT_ITINERAIRE",
     "mea_reservations.png",
     "UTILISATEUR (1) ---< RESERVATION (N)  |  RESERVATION (1) --- PAIEMENT (1)  |  RESERVATION (1) ---< RESERVATION_ACTIVITES (N)\nITINERAIRE (1) ---< ELEMENT_ITINERAIRE (N)  — lien polymorphe vers transport/hébergement/activité"),
    ("MEA — Utilisateurs, Securite & Complements",
     "UTILISATEUR · SESSION · NOTIFICATION · PANIER · FAVORI · AVIS · VOYAGEUR_SEJOUR",
     "mea_complements.png",
     "UTILISATEUR (1) --- PANIER (1)  |  PANIER (1) ---< ELEMENT_PANIER (N)  |  UTILISATEUR (1) ---< FAVORI (N)\nUTILISATEUR (1) ---< AVIS (N)  |  UTILISATEUR (1) ---< SESSION (N)  |  UTILISATEUR (1) ---< NOTIFICATION (N)"),
]
for title, sub, png, note in mea_slides:
    s = new_slide(); bg(s)
    section_tag(s, "03 — MODELE ENTITE-ASSOCIATION")
    slide_title(s, title, sub)
    hline(s, Inches(1.55))

    # Image Mermaid centrée, haute résolution
    img(s, png, Inches(0.4), Inches(1.72), h=Inches(4.8))

    # Note cardinalities en bas
    rect(s, Inches(0.35), Inches(6.65), Inches(12.7), Inches(0.72), CARD_BG)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.4), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    lines = note.split('\n')
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        r.font.size = Pt(9.5); r.font.color.rgb = LIGHT_GREY

    pagenum(s)

# ── SLIDES 13-15: SCHEMA RELATIONNEL ────────────────────────────────────────
sr_slides = [
    ("Schema Relationnel — Catalogue & Offres",
     "Tables : destinations · transports · hebergements · activites",
     "sr_catalogue.png",
     "FK destination_id : transports, hebergements, activites referencent tous destinations(id) avec ON DELETE CASCADE\nContrainte UNIQUE sur destinations.slug — PRIMARY KEY INT AUTO_INCREMENT sur toutes les tables"),
    ("Schema Relationnel — Reservations & Paiements",
     "Tables : utilisateurs · reservations · reservation_activites · paiements · itineraires · itineraire_items",
     "sr_reservations.png",
     "reference VARCHAR(20) UNIQUE — format VV-XXXXXXX generé côté PHP  |  Transaction ACID lors de chaque reservation (BEGIN → COMMIT/ROLLBACK)\nreservation_activites : table de jonction N:M avec nb_places  |  itineraire_items : elements ordonnés (champ ordre INT)"),
    ("Schema Relationnel — Complements & Securite",
     "Tables : notifications · sessions · paniers · elements_panier · favoris · avis · voyageurs_sejour",
     "sr_complements.png",
     "sessions.token VARCHAR(255) UNIQUE — expiration datetime  |  avis.note CHECK(note BETWEEN 1 AND 5)\nelements_panier & favoris : lien polymorphe via type_element ENUM + element_ref_id INT (sans FK stricte)"),
]
for title, sub, png, note in sr_slides:
    s = new_slide(); bg(s)
    section_tag(s, "04 — SCHEMA RELATIONNEL")
    slide_title(s, title, sub)
    hline(s, Inches(1.55))

    img(s, png, Inches(0.4), Inches(1.72), h=Inches(4.8))

    rect(s, Inches(0.35), Inches(6.65), Inches(12.7), Inches(0.72), CARD_BG)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(12.4), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    lines = note.split('\n')
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = line
        r.font.size = Pt(9.5); r.font.color.rgb = LIGHT_GREY

    pagenum(s)

# ── SLIDE 16: ARCHITECTURE — Vue globale (diagramme) ─────────────────────────
s = new_slide(); bg(s)
section_tag(s, "05 — ARCHITECTURE SYSTEME")
slide_title(s, "Architecture Generale — Vue d'ensemble",
            "Frontend React · API PHP REST · MySQL · MAMP / Apache")
hline(s, Inches(1.55))

img(s, "archi_global.png", Inches(0.4), Inches(1.72), h=Inches(4.2))

# Tech stack note
rect(s, Inches(0.35), Inches(6.1), Inches(12.7), Inches(1.2), CARD_BG)
techs = [
    ("Frontend :", "React 18 · Vite 5 · React Router v6 · Axios · CSS Variables (dark/light) · i18n FR/EN"),
    ("Backend  :", "PHP 8 · Architecture MVC · PDO Prepared Statements · Bcrypt · Apache / MAMP :80"),
    ("BDD       :", "MySQL 8 · utf8mb4 · 14 tables · Transactions ACID · Migrations versionnées"),
]
for i, (k, v) in enumerate(techs):
    y = Inches(6.17) + i * Inches(0.36)
    txt(s, k, Inches(0.55), y, Inches(1.2), Inches(0.32),
        size=10, bold=True, color=ACCENT1)
    txt(s, v, Inches(1.75), y, Inches(11.2), Inches(0.32),
        size=10, color=LIGHT_GREY)

pagenum(s)

# ── SLIDE 17: ARCHITECTURE — Flux sequence (diagramme) ───────────────────────
s = new_slide(); bg(s)
section_tag(s, "05 — ARCHITECTURE SYSTEME")
slide_title(s, "Architecture — Flux de Reservation Complet",
            "Sequence de bout en bout : Recherche → Paiement → Confirmation")
hline(s, Inches(1.55))

img(s, "archi_flux.png", Inches(0.5), Inches(1.72), h=Inches(5.5))

pagenum(s)

# ── SLIDE 18: ARCHITECTURE — Composants (diagramme) ─────────────────────────
s = new_slide(); bg(s)
section_tag(s, "05 — ARCHITECTURE SYSTEME")
slide_title(s, "Architecture — Composants Frontend",
            "9 Ecrans · 8 Composants reutilisables · 9 Services Axios")
hline(s, Inches(1.55))

img(s, "archi_composants.png", Inches(0.3), Inches(1.72), w=Inches(12.7))

pagenum(s)

# ── SLIDE 19: ARCHITECTURE — Roles (diagramme) ───────────────────────────────
s = new_slide(); bg(s)
section_tag(s, "05 — ARCHITECTURE SYSTEME")
slide_title(s, "Roles Utilisateurs & Securite",
            "3 roles distincts · Authentification bcrypt · Gestion des droits")
hline(s, Inches(1.55))

img(s, "archi_roles.png", Inches(0.5), Inches(1.72), h=Inches(5.45))

pagenum(s)

# ── HELPERS pour specs ────────────────────────────────────────────────────────
def spec_card(slide, ref, title, desc, color, x, y, w=Inches(6.15), h=Inches(1.55)):
    rect(slide, x, y, w, h, CARD_BG)
    rect(slide, x, y, Inches(0.08), h, color)
    txt(slide, ref,   x + Inches(0.2), y + Inches(0.08),
        Inches(1.0), Inches(0.3), size=9, color=color, italic=True)
    txt(slide, title, x + Inches(1.25), y + Inches(0.08),
        w - Inches(1.35), Inches(0.32), size=11, bold=True, color=WHITE)
    txt(slide, desc,  x + Inches(0.2), y + Inches(0.42),
        w - Inches(0.3), h - Inches(0.5), size=9.5, color=LIGHT_GREY)

SPEC_COLORS = [ACCENT1, ACCENT2, ACCENT3, RED_ACC, PURPLE_ACC, GREEN_ACC]

# ── SLIDE 20: SPECS — Perimetre & Roles ──────────────────────────────────────
s = new_slide(); bg(s)
section_tag(s, "06 — SPECIFICATIONS FONCTIONNELLES")
slide_title(s, "Specifications — Perimetre & Roles",
            "Presentation generale de VoyageVista et de ses acteurs")
hline(s, Inches(1.55))

rect(s, Inches(0.35), Inches(1.72), Inches(12.7), Inches(1.45), CARD_BG)
rect(s, Inches(0.35), Inches(1.72), Inches(0.08), Inches(1.45), ACCENT1)
txt(s, "VoyageVista est une plateforme de voyage en ligne permettant aux utilisateurs de "
       "rechercher des destinations, composer des sejours sur-mesure (transport + hebergement + activites), "
       "creer des itineraires personnalises et effectuer des reservations en ligne. "
       "L'interface est disponible en francais et en anglais avec support du mode sombre.",
    Inches(0.58), Inches(1.82), Inches(12.2), Inches(1.25), size=11, color=LIGHT_GREY)

roles3 = [
    ("Visiteur",         ACCENT1,    "Acces public : consultation, recherche, inscription"),
    ("Utilisateur",      ACCENT2,    "Reservation, itineraires, compte, favoris, avis"),
    ("Administrateur",   ACCENT3,    "Dashboard, gestion utilisateurs et donnees"),
]
for i, (role, color, desc) in enumerate(roles3):
    x = Inches(0.35) + i * Inches(4.32)
    rect(s, x, Inches(3.3), Inches(4.1), Inches(1.0), CARD_BG)
    rect(s, x, Inches(3.3), Inches(0.08), Inches(1.0), color)
    txt(s, role, x + Inches(0.2), Inches(3.38), Inches(3.8), Inches(0.4),
        size=13, bold=True, color=color)
    txt(s, desc, x + Inches(0.2), Inches(3.8),  Inches(3.8), Inches(0.42),
        size=10, color=LIGHT_GREY)

rect(s, Inches(0.35), Inches(4.5), Inches(12.7), Inches(2.8), CARD_BG)
rect(s, Inches(0.35), Inches(4.5), Inches(0.08), Inches(2.8), ACCENT3)
txt(s, "Contraintes Techniques",
    Inches(0.58), Inches(4.57), Inches(6), Inches(0.38),
    size=13, bold=True, color=ACCENT3)
constraints = [
    ("Performance",         "Chargement < 3s, pagination, lazy loading images"),
    ("Securite",            "HTTPS, anti-SQLi PDO, hash bcrypt, CORS restreint :3000"),
    ("Compatibilite",       "Navigateurs modernes (Chrome/Firefox/Safari/Edge), responsive"),
    ("Internationalisation","FR / EN integre, extensible via i18n.js"),
    ("Accessibilite",       "Contraste WCAG AA, navigation clavier, ARIA labels"),
    ("Disponibilite",       "Architecture stateless, frontend/backend deploiement independant"),
]
for i, (k, v) in enumerate(constraints):
    col = i % 2; row = i // 2
    x = Inches(0.58) + col * Inches(6.35)
    y = Inches(5.0) + row * Inches(0.55)
    txt(s, k + " : ", x, y, Inches(1.7), Inches(0.45), size=10, bold=True, color=ACCENT1)
    txt(s, v,         x + Inches(1.7), y, Inches(4.55), Inches(0.45), size=10, color=LIGHT_GREY)

pagenum(s)

# ── SLIDE 21: SF-01 Catalogue & Recherche ────────────────────────────────────
s = new_slide(); bg(s)
section_tag(s, "06 — SPECIFICATIONS FONCTIONNELLES")
slide_title(s, "SF-01 · Catalogue & Recherche",
            "Consultation des destinations et fonctionnalites de recherche")
hline(s, Inches(1.55))

feats21 = [
    ("SF-01.1", "Affichage Catalogue",      ACCENT1,
     "Liste paginee des destinations actives : photo, nom, pays, note, prix depuis, tag categorie. Filtrage par categorie (Asie, Europe, Ameriques...)."),
    ("SF-01.2", "Recherche Globale",         ACCENT2,
     "Barre de recherche combinee : destination + dates aller/retour + nb voyageurs. Resultats filtres en temps reel sur ScreenResults."),
    ("SF-01.3", "Filtres Avances",           ACCENT3,
     "Sur ScreenResults : filtres budget, note minimale, type destination, duree sejour. Tri par pertinence / prix / note."),
    ("SF-01.4", "Detail Destination",        RED_ACC,
     "Page detail avec 4 onglets : Vue d'ensemble (description + galerie), Transports, Hebergements, Activites. Pre-selection en un clic."),
    ("SF-01.5", "Recherche Transport",       PURPLE_ACC,
     "ScreenTransport : recherche multimodale (avion, train, bus, voiture) avec filtres type, compagnie, prix min/max, places disponibles."),
    ("SF-01.6", "Internationalisation",      GREEN_ACC,
     "Titres, descriptions et tags en FR et EN. Toggle langue dans le header, persisté en localStorage."),
]
for i, (ref, title, color, desc) in enumerate(feats21):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.52)
    y = Inches(1.75) + row * Inches(1.78)
    spec_card(s, ref, title, desc, color, x, y)

pagenum(s)

# ── SLIDE 22: SF-02 Reservation & Panier ─────────────────────────────────────
s = new_slide(); bg(s)
section_tag(s, "06 — SPECIFICATIONS FONCTIONNELLES")
slide_title(s, "SF-02 · Reservation & Panier",
            "Composition de sejour, gestion du panier et paiement")
hline(s, Inches(1.55))

feats22 = [
    ("SF-02.1", "Composition de Sejour",    ACCENT1,
     "Depuis ScreenDetail, selection transport + hebergement + activites. Chaque element ajoute au panier (stocke en localStorage)."),
    ("SF-02.2", "Gestion du Panier",         ACCENT2,
     "ScreenCart : liste des elements, modification quantite, suppression, calcul total en temps reel. Badge compteur dans le header."),
    ("SF-02.3", "Validation Reservation",    ACCENT3,
     "POST /reservations : verification stocks, validation dates (retour > depart), generation reference unique VV-XXXXXXX."),
    ("SF-02.4", "Gestion Activites",         RED_ACC,
     "Selection multiple d'activites liees a la destination. Enregistrement dans reservation_activites (N:M), decrementation places_restantes."),
    ("SF-02.5", "Paiement",                  PURPLE_ACC,
     "ScreenPayment : choix mode paiement (carte, PayPal, virement, crypto), recapitulatif final, confirmation commande."),
    ("SF-02.6", "Annulation & Rollback",     GREEN_ACC,
     "Annulation depuis espace compte : remise a jour stocks (nb_chambres_dispo, places_dispo, places_restantes) en transaction ACID."),
]
for i, (ref, title, color, desc) in enumerate(feats22):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.52)
    y = Inches(1.75) + row * Inches(1.78)
    spec_card(s, ref, title, desc, color, x, y)

pagenum(s)

# ── SLIDE 23: SF-03 Itineraires & Compte ─────────────────────────────────────
s = new_slide(); bg(s)
section_tag(s, "06 — SPECIFICATIONS FONCTIONNELLES")
slide_title(s, "SF-03 · Itineraires & Espace Compte",
            "Planification personnalisee et gestion du profil utilisateur")
hline(s, Inches(1.55))

feats23 = [
    ("SF-03.1", "Createur d'Itineraire",    ACCENT1,
     "ScreenItinerary : creer un itineraire nomme avec destination et dates. Ajouter des elements ordonnes (transport, hebergement, activite, autre) avec quantite et plage horaire."),
    ("SF-03.2", "Sauvegarde & Statuts",      ACCENT2,
     "Statuts : brouillon (sauvegarde partielle), confirme, annule. Calcul automatique cout total. API PUT /itineraires pour mise a jour."),
    ("SF-03.3", "Gestion du Profil",         ACCENT3,
     "ScreenAccount : modification nom, prenom, email, telephone, date naissance. Changement mot de passe avec verification ancien."),
    ("SF-03.4", "Historique Reservations",   RED_ACC,
     "Liste des reservations avec statut, montant, reference, dates. Detail accessible. Annulation possible si statut en_attente."),
    ("SF-03.5", "Notifications",             PURPLE_ACC,
     "Temps reel : confirmation reservation, rappels, promos, alertes. Badge non-lu dans header. Marquage lu individuel ou global."),
    ("SF-03.6", "Favoris & Avis",            GREEN_ACC,
     "Ajout en favoris : destinations, transports, hebergements, activites. Avis notes 1-5 etoiles avec commentaire sur tout element."),
]
for i, (ref, title, color, desc) in enumerate(feats23):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.52)
    y = Inches(1.75) + row * Inches(1.78)
    spec_card(s, ref, title, desc, color, x, y)

pagenum(s)

# ── SLIDE 24: SF-04 Administration ───────────────────────────────────────────
s = new_slide(); bg(s)
section_tag(s, "06 — SPECIFICATIONS FONCTIONNELLES")
slide_title(s, "SF-04 · Administration",
            "Dashboard admin — Gestion des utilisateurs et des donnees")
hline(s, Inches(1.55))

feats24 = [
    ("SF-04.1", "Dashboard Admin",           ACCENT1,
     "ScreenAdmin accessible uniquement si role === 'admin'. Vue d'ensemble : utilisateurs, reservations, statistiques generales."),
    ("SF-04.2", "Gestion Utilisateurs",       ACCENT2,
     "Liste complete avec nom, email, role, statut actif. Changement de role (user <-> admin). Suppression de compte via DELETE /admin."),
    ("SF-04.3", "Gestion Destinations",       ACCENT3,
     "Création (POST /destinations), modification (PUT) et suppression (DELETE). Champ actif pour masquer sans supprimer."),
    ("SF-04.4", "Gestion Offres",             RED_ACC,
     "CRUD complet sur transports, hebergements, activites via endpoints dedies. Controle des stocks et disponibilites."),
]
for i, (ref, title, color, desc) in enumerate(feats24):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.52)
    y = Inches(1.75) + row * Inches(2.1)
    spec_card(s, ref, title, desc, color, x, y, h=Inches(1.9))

# Techno
rect(s, Inches(0.35), Inches(6.1), Inches(12.7), Inches(1.2), CARD_BG)
rect(s, Inches(0.35), Inches(6.1), Inches(0.08), Inches(1.2), PURPLE_ACC)
txt(s, "Technologies", Inches(0.58), Inches(6.16), Inches(2.5), Inches(0.35),
    size=12, bold=True, color=PURPLE_ACC)
techrows = [
    ("Frontend :", "React 18, Vite 5, React Router v6, Axios, CSS Variables"),
    ("Backend  :", "PHP 8, Architecture MVC, PDO, Bcrypt, Apache / MAMP"),
    ("BDD       :", "MySQL 8, utf8mb4, Transactions ACID, 14 tables, Migrations"),
]
for i, (k, v) in enumerate(techrows):
    txt(s, k, Inches(0.58), Inches(6.55) + i * Inches(0.38), Inches(1.2), Inches(0.35),
        size=9.5, bold=True, color=ACCENT1)
    txt(s, v, Inches(1.78), Inches(6.55) + i * Inches(0.38), Inches(11.1), Inches(0.35),
        size=9.5, color=LIGHT_GREY)

pagenum(s)

# ── SLIDE 25: Synthese ────────────────────────────────────────────────────────
s = new_slide(); bg(s)
section_tag(s, "06 — SPECIFICATIONS FONCTIONNELLES")
slide_title(s, "Synthese des Fonctionnalites",
            "Vue consolidee — correspondance fonctionnalites / modules techniques")
hline(s, Inches(1.55))

col_specs = [(0.6, "Ref."), (1.8, "Fonctionnalite"), (3.7, "Ecran(s) React"), (4.2, "Table(s) MySQL")]
x0 = Inches(0.35)
y_hdr = Inches(1.72)
x = x0
for w_in, label in col_specs:
    rect(s, x, y_hdr, Inches(w_in), Inches(0.4), ACCENT1)
    txt(s, label, x + Inches(0.08), y_hdr + Inches(0.04),
        Inches(w_in - 0.12), Inches(0.32), size=10, bold=True, color=TEXT_DARK)
    x += Inches(w_in)

rows = [
    ("SF-01",   "Catalogue & Recherche",    "ScreenHome, ScreenResults,\nScreenDetail, ScreenTransport",          "destinations, transports,\nhebergements, activites",        ACCENT1),
    ("SF-02",   "Reservation & Panier",     "ScreenCart, ScreenPayment",                                           "reservations, reservation_activites,\npaiements",            ACCENT2),
    ("SF-03.1", "Itineraires",              "ScreenItinerary",                                                     "itineraires, itineraire_items",                               ACCENT3),
    ("SF-03.3", "Espace Compte",            "ScreenAccount",                                                       "utilisateurs, sessions,\nnotifications",                      RED_ACC),
    ("SF-03.6", "Favoris & Avis",           "ScreenAccount (onglets)",                                             "favoris, avis",                                               PURPLE_ACC),
    ("SF-04",   "Administration",           "ScreenAdmin",                                                         "utilisateurs (CRUD admin)",                                   GREEN_ACC),
]
for i, (ref, feat, screens, tables, color) in enumerate(rows):
    y_row = Inches(2.18) + i * Inches(0.82)
    bg_row = CARD_BG if i % 2 == 0 else RGBColor(0x12, 0x20, 0x30)
    x = x0
    for j, (w_in, _) in enumerate(col_specs):
        vals = [ref, feat, screens, tables]
        rect(s, x, y_row, Inches(w_in), Inches(0.76), bg_row)
        c = color if j == 0 else (WHITE if j == 1 else LIGHT_GREY)
        txt(s, vals[j], x + Inches(0.08), y_row + Inches(0.06),
            Inches(w_in - 0.12), Inches(0.66),
            size=9 if j > 1 else 10, bold=(j == 1), color=c)
        x += Inches(w_in)

pagenum(s)

# ── SAVE ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\wanga\Downloads\VoyageVista_PowerPoint_Preliminaire_v2.pptx"
prs.save(out)
print("SAVED:", out)
print("Total slides:", len(prs.slides))
