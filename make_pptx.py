from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Palette ───────────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # bleu marine très foncé
CARD_BG   = RGBColor(0x16, 0x2A, 0x3E)   # bleu marine card
ACCENT1   = RGBColor(0x00, 0xB4, 0xD8)   # bleu cyan vif
ACCENT2   = RGBColor(0x48, 0xCA, 0xB2)   # turquoise
ACCENT3   = RGBColor(0xFF, 0xA5, 0x00)   # orange accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY= RGBColor(0xC8, 0xD6, 0xE5)
MID_GREY  = RGBColor(0x5A, 0x7A, 0x9A)
TEXT_DARK = RGBColor(0x0D, 0x1B, 0x2A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # blank layout

slide_counter = [0]

# ─── Helpers ───────────────────────────────────────────────────────────────────
def new_slide():
    slide_counter[0] += 1
    return prs.slides.add_slide(BLANK)

def fill_slide_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_textbox(slide, text, l, t, w, h,
                size=16, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_para(tf, text, size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def page_num(slide, n):
    add_textbox(slide, str(n),
                W - Inches(0.55), H - Inches(0.38),
                Inches(0.4), Inches(0.28),
                size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)

def section_tag(slide, label, color=ACCENT1):
    add_rect(slide, Inches(0.4), Inches(0.22), Inches(0.08), Inches(0.36), color)
    add_textbox(slide, label,
                Inches(0.55), Inches(0.2), Inches(5), Inches(0.4),
                size=11, bold=True, color=color)

def slide_title(slide, title, subtitle=None):
    add_textbox(slide, title,
                Inches(0.55), Inches(0.62), Inches(12), Inches(0.7),
                size=30, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.55), Inches(1.3), Inches(12), Inches(0.4),
                    size=14, color=LIGHT_GREY, italic=True)

def h_line(slide, y, color=ACCENT1, thickness=2):
    from pptx.util import Pt as Pt2
    ln = slide.shapes.add_shape(1,
        Inches(0.4), y, W - Inches(0.8), Pt2(thickness))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()

# ─── SLIDE 1 : PAGE DE GARDE ───────────────────────────────────────────────────
s = new_slide()
fill_slide_bg(s, DARK_BG)

# gradient band gauche
add_rect(s, 0, 0, Inches(4.5), H, CARD_BG)
add_rect(s, 0, 0, Inches(0.12), H, ACCENT1)

add_textbox(s, "VOYAGEVISTA", Inches(0.55), Inches(1.4), Inches(3.8), Inches(0.8),
            size=28, bold=True, color=ACCENT1, align=PP_ALIGN.LEFT)
add_textbox(s, "Plateforme de voyage en ligne",
            Inches(0.55), Inches(2.22), Inches(3.8), Inches(0.5),
            size=14, color=LIGHT_GREY, italic=True)

add_rect(s, Inches(0.55), Inches(2.85), Inches(3.4), Inches(0.04), ACCENT2)

# Zone à compléter
add_rect(s, Inches(0.55), Inches(3.1), Inches(3.6), Inches(2.2), RGBColor(0x0A, 0x15, 0x22))
add_textbox(s, "[ Noms des membres ]",
            Inches(0.7), Inches(3.2), Inches(3.3), Inches(0.4),
            size=13, color=MID_GREY, italic=True)
add_textbox(s, "[ Groupe / Promotion ]",
            Inches(0.7), Inches(3.65), Inches(3.3), Inches(0.4),
            size=13, color=MID_GREY, italic=True)
add_textbox(s, "[ Date de rendu ]",
            Inches(0.7), Inches(4.1), Inches(3.3), Inches(0.4),
            size=13, color=MID_GREY, italic=True)
add_textbox(s, "[ Enseignant ]",
            Inches(0.7), Inches(4.55), Inches(3.3), Inches(0.4),
            size=13, color=MID_GREY, italic=True)

# Droite : sommaire rapide
add_textbox(s, "PowerPoint Préliminaire",
            Inches(5.0), Inches(1.2), Inches(7.8), Inches(0.6),
            size=22, bold=True, color=WHITE)
add_textbox(s, "Livrable de conception — Projet Web 2026",
            Inches(5.0), Inches(1.85), Inches(7.8), Inches(0.4),
            size=12, color=LIGHT_GREY, italic=True)

sections = [
    ("01", "Wireframes & Interfaces utilisateur"),
    ("02", "Storyboard & Parcours utilisateurs"),
    ("03", "Modèle Entité-Association"),
    ("04", "Schéma Relationnel"),
    ("05", "Architecture Système"),
    ("06", "Spécifications Fonctionnelles"),
]
for i, (num, label) in enumerate(sections):
    y = Inches(2.5) + i * Inches(0.62)
    add_rect(s, Inches(5.0), y, Inches(0.5), Inches(0.42), CARD_BG)
    add_textbox(s, num, Inches(5.02), y + Inches(0.04), Inches(0.46), Inches(0.38),
                size=13, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    add_textbox(s, label, Inches(5.6), y + Inches(0.04), Inches(7.1), Inches(0.38),
                size=13, color=WHITE)

page_num(s, 1)

# ─── SLIDE 2 : SOMMAIRE ────────────────────────────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "SOMMAIRE")
slide_title(s, "Table des matières")
h_line(s, Inches(1.45))

items = [
    ("01 — Wireframes & Interfaces utilisateur",  "Slides 3–5",   ACCENT1),
    ("02 — Storyboard & Parcours utilisateurs",   "Slides 6–8",   ACCENT2),
    ("03 — Modèle Entité-Association",            "Slides 9–12",  ACCENT3),
    ("04 — Schéma Relationnel",                   "Slides 13–15", RGBColor(0xFF,0x69,0x69)),
    ("05 — Architecture Système",                 "Slides 16–19", RGBColor(0x9B,0x59,0xB6)),
    ("06 — Spécifications Fonctionnelles",        "Slides 20–26", RGBColor(0x2E,0xCC,0x71)),
]
for i, (label, pages, color) in enumerate(items):
    col = 0 if i < 3 else 1
    row = i % 3
    x = Inches(0.5) + col * Inches(6.5)
    y = Inches(1.7) + row * Inches(1.55)
    add_rect(s, x, y, Inches(6.0), Inches(1.35), CARD_BG)
    add_rect(s, x, y, Inches(0.1), Inches(1.35), color)
    add_textbox(s, label, x + Inches(0.25), y + Inches(0.15),
                Inches(5.5), Inches(0.6), size=14, bold=True, color=WHITE)
    add_textbox(s, pages, x + Inches(0.25), y + Inches(0.75),
                Inches(5.5), Inches(0.4), size=12, color=color, italic=True)

page_num(s, 2)

# ─── SLIDES 3–5 : WIREFRAMES (à compléter) ────────────────────────────────────
wf_labels = [
    ("Wireframes — Pages principales",
     "Accueil · Résultats · Détail destination",
     ["Page Accueil (Hero + Recherche + Destinations)", "Page Résultats (Filtres + Cartes)", "Page Détail Destination (Onglets)"]),
    ("Wireframes — Espace utilisateur",
     "Panier · Paiement · Compte · Itinéraire",
     ["Page Panier / Checkout", "Page Paiement & Confirmation", "Espace Compte Utilisateur", "Créateur d'Itinéraire"]),
    ("Wireframes — Administration & Recherche transport",
     "Transport · Dashboard Admin",
     ["Page Recherche Transport Multimodal", "Dashboard Admin (Gestion utilisateurs & données)"]),
]
for i, (title, sub, items) in enumerate(wf_labels):
    s = new_slide()
    fill_slide_bg(s)
    section_tag(s, "01 — WIREFRAMES")
    slide_title(s, title, sub)
    h_line(s, Inches(1.52))

    # Zone placeholder wireframe
    add_rect(s, Inches(0.4), Inches(1.7), Inches(8.4), Inches(5.4), CARD_BG)
    add_rect(s, Inches(0.4), Inches(1.7), Inches(8.4), Inches(0.04), ACCENT1)
    add_textbox(s, "[ Insérer wireframe ici ]",
                Inches(0.4), Inches(4.0), Inches(8.4), Inches(0.8),
                size=16, color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)

    # Liste pages couvertes
    add_rect(s, Inches(9.0), Inches(1.7), Inches(3.9), Inches(5.4), CARD_BG)
    add_rect(s, Inches(9.0), Inches(1.7), Inches(0.08), Inches(5.4), ACCENT1)
    add_textbox(s, "Pages représentées",
                Inches(9.15), Inches(1.82), Inches(3.6), Inches(0.4),
                size=13, bold=True, color=ACCENT1)
    txb = slide.shapes if False else None
    tb = s.shapes.add_textbox(Inches(9.15), Inches(2.3), Inches(3.55), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "▸  " + item
        r.font.size = Pt(12); r.font.color.rgb = LIGHT_GREY

    page_num(s, 3 + i)

# ─── SLIDES 6–8 : STORYBOARD (à compléter) ────────────────────────────────────
sb_data = [
    ("Storyboard — Parcours Réservation",
     "De la recherche à la confirmation de réservation",
     ["Étape 1 : Accueil & Recherche", "Étape 2 : Consultation des résultats", "Étape 3 : Détail destination & sélection", "Étape 4 : Ajout au panier"]),
    ("Storyboard — Finalisation & Compte",
     "Paiement, confirmation, espace personnel",
     ["Étape 5 : Récapitulatif panier", "Étape 6 : Paiement & confirmation", "Étape 7 : Réception notification", "Étape 8 : Consultation espace compte"]),
    ("Storyboard — Création d'itinéraire",
     "Fonctionnalité de planification personnalisée",
     ["Étape 1 : Accès au créateur d'itinéraire", "Étape 2 : Ajout d'étapes & composants", "Étape 3 : Sauvegarde & consultation", "Étape 4 : Partage / export"]),
]
for i, (title, sub, steps) in enumerate(sb_data):
    s = new_slide()
    fill_slide_bg(s)
    section_tag(s, "02 — STORYBOARD")
    slide_title(s, title, sub)
    h_line(s, Inches(1.52))

    add_rect(s, Inches(0.4), Inches(1.7), Inches(8.4), Inches(5.4), CARD_BG)
    add_rect(s, Inches(0.4), Inches(1.7), Inches(8.4), Inches(0.04), ACCENT2)
    add_textbox(s, "[ Insérer storyboard / captures ici ]",
                Inches(0.4), Inches(4.0), Inches(8.4), Inches(0.8),
                size=16, color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)

    add_rect(s, Inches(9.0), Inches(1.7), Inches(3.9), Inches(5.4), CARD_BG)
    add_rect(s, Inches(9.0), Inches(1.7), Inches(0.08), Inches(5.4), ACCENT2)
    add_textbox(s, "Étapes du parcours",
                Inches(9.15), Inches(1.82), Inches(3.6), Inches(0.4),
                size=13, bold=True, color=ACCENT2)
    tb = s.shapes.add_textbox(Inches(9.15), Inches(2.3), Inches(3.55), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for step in steps:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "▸  " + step
        r.font.size = Pt(12); r.font.color.rgb = LIGHT_GREY

    page_num(s, 6 + i)

# ─── SLIDE 9 : MEA — Vue d'ensemble ───────────────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "03 — MODÈLE ENTITÉ-ASSOCIATION")
slide_title(s, "MEA — Vue d'ensemble des entités",
            "6 modules fonctionnels · 14 entités · Relations identifiées")
h_line(s, Inches(1.52))

modules = [
    ("Catalogue & Offres",   ["DESTINATION", "TRANSPORT", "HÉBERGEMENT", "ACTIVITÉ"],  ACCENT1,  Inches(0.4),  Inches(1.7)),
    ("Réservations",         ["RÉSERVATION", "PAIEMENT", "ÉLÉMENT_ITINÉRAIRE"],        ACCENT2,  Inches(4.55), Inches(1.7)),
    ("Itinéraires",          ["ITINÉRAIRE"],                                            ACCENT3,  Inches(8.7),  Inches(1.7)),
    ("Utilisateurs",         ["UTILISATEUR", "SESSION", "NOTIFICATION"],               RGBColor(0xFF,0x69,0x69), Inches(0.4),  Inches(4.3)),
    ("Compléments",          ["PANIER", "ÉLÉMENT_PANIER", "FAVORI", "AVIS"],           RGBColor(0x9B,0x59,0xB6), Inches(4.55), Inches(4.3)),
    ("Voyageur",             ["VOYAGEUR_SEJOUR"],                                       RGBColor(0x2E,0xCC,0x71), Inches(8.7),  Inches(4.3)),
]
for name, entities, color, x, y in modules:
    bh = Inches(0.45) + len(entities) * Inches(0.38)
    add_rect(s, x, y, Inches(3.9), bh, CARD_BG)
    add_rect(s, x, y, Inches(0.08), bh, color)
    add_textbox(s, name, x + Inches(0.18), y + Inches(0.05),
                Inches(3.6), Inches(0.36), size=12, bold=True, color=color)
    for j, ent in enumerate(entities):
        add_textbox(s, "  " + ent,
                    x + Inches(0.18), y + Inches(0.4) + j * Inches(0.38),
                    Inches(3.6), Inches(0.35), size=11, color=LIGHT_GREY)

page_num(s, 9)

# ─── SLIDE 10 : MEA — Bloc Catalogue ──────────────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "03 — MODÈLE ENTITÉ-ASSOCIATION")
slide_title(s, "MEA — Bloc Catalogue & Offres",
            "Entités : DESTINATION · TRANSPORT · HÉBERGEMENT · ACTIVITÉ")
h_line(s, Inches(1.52))

def entity_box(slide, title, attrs, x, y, w=Inches(3.0), color=ACCENT1):
    row_h = Inches(0.3)
    h = Inches(0.5) + len(attrs) * row_h
    add_rect(slide, x, y, w, h, CARD_BG)
    add_rect(slide, x, y, w, Inches(0.42), color)
    add_textbox(slide, title, x + Inches(0.1), y + Inches(0.04),
                w - Inches(0.2), Inches(0.36),
                size=12, bold=True, color=TEXT_DARK if color==ACCENT1 else WHITE,
                align=PP_ALIGN.CENTER)
    for i, (attr, typ, pk) in enumerate(attrs):
        bg = RGBColor(0x11,0x22,0x33) if i % 2 == 0 else CARD_BG
        ay = y + Inches(0.46) + i * row_h
        add_rect(slide, x, ay, w, row_h, bg)
        prefix = "🔑 " if pk == "PK" else ("🔗 " if pk == "FK" else "   ")
        add_textbox(slide, prefix + attr,
                    x + Inches(0.08), ay + Inches(0.02),
                    w * 0.62, row_h - Inches(0.04),
                    size=9, color=WHITE if pk != "PK" else ACCENT3)
        add_textbox(slide, typ,
                    x + w * 0.62 + Inches(0.05), ay + Inches(0.02),
                    w * 0.36, row_h - Inches(0.04),
                    size=8, color=MID_GREY, italic=True)
    return h

dest_attrs = [
    ("id_destination","INT","PK"),("slug","VARCHAR(100)",""),
    ("nom","VARCHAR(100)",""),("pays","VARCHAR(80)",""),
    ("region","VARCHAR(80)",""),("description","TEXT",""),
    ("photo_principale","VARCHAR(255)",""),("categorie","ENUM(...)",""),
    ("note_moyenne","DECIMAL(3,2)",""),("prix_moyen","DECIMAL(10,2)",""),
    ("actif","BOOLEAN",""),
]
h = entity_box(s, "DESTINATION", dest_attrs, Inches(0.3), Inches(1.65), Inches(3.4), ACCENT1)

transp_attrs = [
    ("id_transport","INT","PK"),("type","ENUM(avion,train,bus,voiture)",""),
    ("compagnie","VARCHAR(100)",""),("origine","VARCHAR(100)",""),
    ("destination_finale","VARCHAR(100)",""),("date_depart","DATETIME",""),
    ("date_arrivee","DATETIME",""),("prix_par_personne","DECIMAL(10,2)",""),
    ("places_disponibles","INT",""),("places_totales","INT",""),
    ("id_destination","INT","FK"),
]
entity_box(s, "TRANSPORT", transp_attrs, Inches(3.9), Inches(1.65), Inches(3.4), ACCENT2)

heb_attrs = [
    ("id_hebergement","INT","PK"),("nom","VARCHAR(150)",""),
    ("type","ENUM(hotel,villa,auberge,appart)",""),
    ("description","TEXT",""),("adresse","VARCHAR(255)",""),
    ("etoiles","INT",""),("prix_nuit","DECIMAL(10,2)",""),
    ("connexion_wifi","BOOLEAN",""),("note_moyenne","DECIMAL(3,2)",""),
    ("annulation_gratuite","BOOLEAN",""),("id_destination","INT","FK"),
]
entity_box(s, "HÉBERGEMENT", heb_attrs, Inches(7.5), Inches(1.65), Inches(3.3), ACCENT3)

act_attrs = [
    ("id_activite","INT","PK"),("nom","VARCHAR(150)",""),
    ("description","TEXT",""),("categorie","ENUM(...)",""),
    ("duree_heures","DECIMAL(4,1)",""),("prix_par_personne","DECIMAL(10,2)",""),
    ("capacite_max","INT",""),("date_heure","DATETIME",""),
    ("places_restantes","INT",""),("id_destination","INT","FK"),
]
entity_box(s, "ACTIVITÉ", act_attrs, Inches(11.0), Inches(1.65), Inches(2.15), RGBColor(0xFF,0x69,0x69))

# Relations
def arrow(slide, x1, y1, x2, y2, label="", color=MID_GREY):
    from pptx.util import Pt as Pt_
    conn = slide.shapes.add_connector(1,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt_(1.5)
    if label:
        mx, my = (x1+x2)/2, (y1+y2)/2
        add_textbox(slide, label, Inches(mx-0.35), Inches(my-0.15),
                    Inches(0.7), Inches(0.28), size=8, color=ACCENT1, italic=True)

arrow(s, 3.7, 3.0, 3.9, 3.0, "1,N")
arrow(s, 7.3, 3.0, 7.5, 3.0, "1,N")
arrow(s, 10.8, 3.0, 11.0, 3.0, "1,N")

# Cardinalities note
add_rect(s, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.65), CARD_BG)
add_textbox(s,
    "Cardinalités : DESTINATION (1) ──< TRANSPORT (N)  ·  DESTINATION (1) ──< HÉBERGEMENT (N)  ·  DESTINATION (1) ──< ACTIVITÉ (N)",
    Inches(0.4), Inches(6.62), Inches(12.5), Inches(0.55),
    size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

page_num(s, 10)

# ─── SLIDE 11 : MEA — Bloc Réservations & Itinéraires ─────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "03 — MODÈLE ENTITÉ-ASSOCIATION")
slide_title(s, "MEA — Réservations & Itinéraires",
            "Entités : RÉSERVATION · PAIEMENT · ITINÉRAIRE · ÉLÉMENT_ITINÉRAIRE")
h_line(s, Inches(1.52))

res_attrs = [
    ("id_reservation","INT","PK"),("reference","VARCHAR(20) UNIQUE",""),
    ("statut","ENUM(en_attente,confirmée,annulée,terminée)",""),
    ("montant_total","DECIMAL(10,2)",""),("date_reservation","DATETIME",""),
    ("date_paiement","DATETIME",""),("id_utilisateur","INT","FK"),
    ("id_destination","INT","FK"),("id_hebergement","INT","FK"),
    ("id_transport","INT","FK"),("date_depart","DATE",""),
    ("date_retour","DATE",""),("nb_voyageurs","INT",""),
]
entity_box(s, "RÉSERVATION", res_attrs, Inches(0.3), Inches(1.7), Inches(3.7), ACCENT1)

pay_attrs = [
    ("id_paiement","INT","PK"),("mode","ENUM(carte,virement,paypal,crypto)",""),
    ("montant","DECIMAL(10,2)",""),("statut","ENUM(en_attente,validé,refusé,remboursé)",""),
    ("date_paiement","DATETIME",""),("reference_externe","VARCHAR(100)",""),
    ("id_reservation","INT","FK"),
]
entity_box(s, "PAIEMENT", pay_attrs, Inches(4.3), Inches(1.7), Inches(3.3), ACCENT2)

itin_attrs = [
    ("id_itineraire","INT","PK"),("titre","VARCHAR(200)",""),
    ("statut","ENUM(brouillon,confirmé,annulé)",""),
    ("date_debut","DATE",""),("date_fin","DATE",""),
    ("nb_voyageurs","INT",""),("cout_total","DECIMAL(10,2)",""),
    ("date_creation","DATETIME",""),("id_utilisateur","INT","FK"),
    ("id_destination","INT","FK"),
]
entity_box(s, "ITINÉRAIRE", itin_attrs, Inches(7.9), Inches(1.7), Inches(3.0), ACCENT3)

elem_attrs = [
    ("id_element","INT","PK"),("type_element","ENUM(transport,hébergement,activité,autre)",""),
    ("id_element_ref","INT",""),("quantite","INT",""),
    ("prix_unitaire","DECIMAL(10,2)",""),("date_debut","DATETIME",""),
    ("date_fin","DATETIME",""),("ordre","INT",""),
    ("id_itineraire","INT","FK"),
]
entity_box(s, "ÉLÉMENT_ITINÉRAIRE", elem_attrs, Inches(11.1), Inches(1.7), Inches(2.1), RGBColor(0xFF,0x69,0x69))

arrow(s, 4.0, 3.2, 4.3, 3.2, "1,1")
arrow(s, 7.6, 3.2, 7.9, 3.2, "1,N")
arrow(s, 11.0, 3.2, 11.1, 3.2, "1,N")

add_rect(s, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.7), CARD_BG)
add_textbox(s,
    "RÉSERVATION (1) ──< PAIEMENT (1)  ·  UTILISATEUR (1) ──< RÉSERVATION (N)  ·  ITINÉRAIRE (1) ──< ÉLÉMENT_ITINÉRAIRE (N)",
    Inches(0.4), Inches(6.57), Inches(12.5), Inches(0.6),
    size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

page_num(s, 11)

# ─── SLIDE 12 : MEA — Utilisateurs & Compléments ─────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "03 — MODÈLE ENTITÉ-ASSOCIATION")
slide_title(s, "MEA — Utilisateurs, Sécurité & Compléments",
            "UTILISATEUR · SESSION · NOTIFICATION · PANIER · FAVORI · AVIS · VOYAGEUR_SEJOUR")
h_line(s, Inches(1.52))

user_attrs = [
    ("id_utilisateur","INT","PK"),("nom","VARCHAR(100)",""),
    ("prenom","VARCHAR(100)",""),("email","VARCHAR(150) UNIQUE",""),
    ("mot_de_passe_hash","VARCHAR(255)",""),("role","ENUM(user,admin,gestionnaire)",""),
    ("date_inscription","DATETIME",""),("photo_profil","VARCHAR(255)",""),
    ("actif","BOOLEAN",""),
]
entity_box(s, "UTILISATEUR", user_attrs, Inches(0.3), Inches(1.7), Inches(3.2), ACCENT1)

sess_attrs = [
    ("id_session","INT","PK"),("token","VARCHAR(255) UNIQUE",""),
    ("expire_le","DATETIME",""),("ip_adresse","VARCHAR(45)",""),
    ("id_utilisateur","INT","FK"),
]
entity_box(s, "SESSION", sess_attrs, Inches(3.7), Inches(1.7), Inches(2.8), ACCENT2)

notif_attrs = [
    ("id_notif","INT","PK"),("type","ENUM(confirmation,rappel,promo,alerte)",""),
    ("message","TEXT",""),("lu","BOOLEAN",""),
    ("date_envoi","DATETIME",""),("id_utilisateur","INT","FK"),
]
entity_box(s, "NOTIFICATION", notif_attrs, Inches(6.7), Inches(1.7), Inches(2.9), ACCENT3)

panier_attrs = [
    ("id_panier","INT","PK"),("date_mise_a_jour","DATETIME",""),
    ("id_utilisateur","INT","FK"),
]
entity_box(s, "PANIER", panier_attrs, Inches(0.3), Inches(4.5), Inches(2.6), RGBColor(0xFF,0x69,0x69))

elem_p_attrs = [
    ("id_element_panier","INT","PK"),("type_element","ENUM(transport,hébergement,activité)",""),
    ("id_element_ref","INT",""),("quantite","INT",""),
    ("prix_unitaire","DECIMAL(10,2)",""),("date_ajout","DATETIME",""),
    ("id_panier","INT","FK"),
]
entity_box(s, "ÉLÉMENT_PANIER", elem_p_attrs, Inches(3.1), Inches(4.5), Inches(3.3), RGBColor(0x9B,0x59,0xB6))

favori_attrs = [
    ("id_favori","INT","PK"),("type_element","ENUM(destination,transport,hébergement,activité)",""),
    ("id_element_ref","INT",""),("date_ajout","DATETIME",""),
    ("id_utilisateur","INT","FK"),
]
entity_box(s, "FAVORI", favori_attrs, Inches(6.6), Inches(4.5), Inches(3.0), RGBColor(0x2E,0xCC,0x71))

avis_attrs = [
    ("id_avis","INT","PK"),("type_cible","ENUM(destination,transport,hébergement,activité)",""),
    ("id_cible","INT",""),("note","INT",""),
    ("commentaire","TEXT",""),("date_avis","DATETIME",""),
    ("id_utilisateur","INT","FK"),
]
entity_box(s, "AVIS", avis_attrs, Inches(9.8), Inches(1.7), Inches(3.1), RGBColor(0xF3,0x9C,0x12))

vs_attrs = [
    ("id","INT","PK"),("refe_sejour","VARCHAR(50)",""),
    ("date_ajout","DATETIME",""),("id_reservation","INT","FK"),
    ("id_utilisateur","INT","FK"),
]
entity_box(s, "VOYAGEUR_SEJOUR", vs_attrs, Inches(9.8), Inches(4.5), Inches(3.1), MID_GREY)

page_num(s, 12)

# ─── SLIDES 13–15 : SCHÉMA RELATIONNEL ────────────────────────────────────────

# Slide 13 : SR Catalogue
s = new_slide()
fill_slide_bg(s)
section_tag(s, "04 — SCHÉMA RELATIONNEL")
slide_title(s, "Schéma Relationnel — Catalogue & Offres",
            "Tables : destinations · transports · hebergements · activites")
h_line(s, Inches(1.52))

def sr_table(slide, name, cols, x, y, w=Inches(3.0), color=ACCENT1):
    row_h = Inches(0.285)
    h = Inches(0.44) + len(cols) * row_h
    add_rect(slide, x, y, w, Inches(0.44), color)
    add_textbox(slide, name, x + Inches(0.1), y + Inches(0.04),
                w - Inches(0.15), Inches(0.36),
                size=12, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    for i, (col, spec) in enumerate(cols):
        bg = RGBColor(0x11,0x22,0x33) if i % 2 == 0 else CARD_BG
        ay = y + Inches(0.44) + i * row_h
        add_rect(slide, x, ay, w, row_h, bg)
        add_textbox(slide, col, x + Inches(0.1), ay + Inches(0.02),
                    w * 0.55, row_h - Inches(0.04), size=9, color=WHITE)
        add_textbox(slide, spec, x + w * 0.55, ay + Inches(0.02),
                    w * 0.43, row_h - Inches(0.04), size=8, color=MID_GREY, italic=True)
    return h

sr_table(s, "destinations",
    [("id","INT PK AUTO_INCREMENT"),("slug","VARCHAR(100) UNIQUE NOT NULL"),
     ("ville","VARCHAR(100)"),("pays_fr","VARCHAR(80)"),
     ("pays_en","VARCHAR(80)"),("type","VARCHAR(50)"),
     ("note","DECIMAL(3,2) DEFAULT 0"),("nb_avis","INT DEFAULT 0"),
     ("duree_jours","INT"),("prix_depuis","DECIMAL(10,2)"),
     ("tag_fr","VARCHAR(100)"),("resume_fr","TEXT"),
     ("image_url","VARCHAR(255)"),("actif","BOOLEAN DEFAULT 1")],
    Inches(0.3), Inches(1.7), Inches(3.25), ACCENT1)

sr_table(s, "transports",
    [("id","INT PK AUTO_INCREMENT"),("destination_id","INT FK → destinations.id"),
     ("type","ENUM(avion,train,bus,voiture)"),
     ("compagnie","VARCHAR(100)"),("depart","VARCHAR(100)"),
     ("arrivee","VARCHAR(100)"),("duree","VARCHAR(30)"),
     ("horaire","VARCHAR(50)"),("prix","DECIMAL(10,2)"),
     ("places_dispo","INT DEFAULT 0")],
    Inches(3.75), Inches(1.7), Inches(3.35), ACCENT2)

sr_table(s, "hebergements",
    [("id","INT PK AUTO_INCREMENT"),("destination_id","INT FK → destinations.id"),
     ("nom","VARCHAR(150)"),("quartier","VARCHAR(100)"),
     ("type","ENUM(hotel,villa,auberge,appartement)"),
     ("prix_nuit","DECIMAL(10,2)"),("nb_etoiles","INT"),
     ("note","DECIMAL(3,2)"),("avantage_fr","TEXT"),
     ("image_url","VARCHAR(255)"),("nb_chambres_dispo","INT DEFAULT 0")],
    Inches(7.3), Inches(1.7), Inches(3.25), ACCENT3)

sr_table(s, "activites",
    [("id","INT PK AUTO_INCREMENT"),("destination_id","INT FK → destinations.id"),
     ("nom_fr","VARCHAR(150)"),("nom_en","VARCHAR(150)"),
     ("categorie","VARCHAR(80)"),("duree","VARCHAR(50)"),
     ("prix","DECIMAL(10,2)"),("description_fr","TEXT"),
     ("description_en","TEXT"),("places_restantes","INT DEFAULT 0")],
    Inches(10.75), Inches(1.7), Inches(2.45), RGBColor(0xFF,0x69,0x69))

# FK arrows labels
for x1, x2, label in [(3.55, 3.75, "FK"), (7.1, 7.3, "FK"), (10.55, 10.75, "FK")]:
    arrow(s, x1, 3.5, x2, 3.5, "destination_id →", ACCENT1)

add_rect(s, Inches(0.3), Inches(6.65), Inches(12.7), Inches(0.6), CARD_BG)
add_textbox(s,
    "Toutes les tables offres (transports, hebergements, activites) possèdent une FK destination_id référençant destinations(id) avec ON DELETE CASCADE",
    Inches(0.4), Inches(6.67), Inches(12.5), Inches(0.5),
    size=9.5, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

page_num(s, 13)

# Slide 14 : SR Réservations
s = new_slide()
fill_slide_bg(s)
section_tag(s, "04 — SCHÉMA RELATIONNEL")
slide_title(s, "Schéma Relationnel — Réservations & Paiements",
            "Tables : utilisateurs · reservations · reservation_activites · paiements · itineraires")
h_line(s, Inches(1.52))

sr_table(s, "utilisateurs",
    [("id","INT PK AUTO_INCREMENT"),("nom","VARCHAR(100) NOT NULL"),
     ("prenom","VARCHAR(100)"),("email","VARCHAR(150) UNIQUE NOT NULL"),
     ("mot_de_passe","VARCHAR(255) NOT NULL"),
     ("role","ENUM(user,admin) DEFAULT user"),
     ("telephone","VARCHAR(20)"),("date_naissance","DATE"),
     ("photo_profil","VARCHAR(255)"),("actif","BOOLEAN DEFAULT 1"),
     ("created_at","DATETIME DEFAULT NOW()")],
    Inches(0.3), Inches(1.7), Inches(3.3), ACCENT1)

sr_table(s, "reservations",
    [("id","INT PK AUTO_INCREMENT"),("reference","VARCHAR(20) UNIQUE NOT NULL"),
     ("utilisateur_id","INT FK → utilisateurs.id"),
     ("destination_id","INT FK → destinations.id"),
     ("hebergement_id","INT FK → hebergements.id"),
     ("transport_id","INT FK → transports.id"),
     ("date_depart","DATE NOT NULL"),("date_retour","DATE NOT NULL"),
     ("nb_voyageurs","INT DEFAULT 1"),
     ("montant_total","DECIMAL(10,2)"),
     ("statut","ENUM(en_attente,confirmee,annulee,terminee)"),
     ("date_reservation","DATETIME DEFAULT NOW()")],
    Inches(3.85), Inches(1.7), Inches(3.4), ACCENT2)

sr_table(s, "reservation_activites",
    [("id","INT PK AUTO_INCREMENT"),("reservation_id","INT FK → reservations.id"),
     ("activite_id","INT FK → activites.id"),
     ("nb_places","INT DEFAULT 1")],
    Inches(7.5), Inches(1.7), Inches(2.8), ACCENT3)

sr_table(s, "paiements",
    [("id","INT PK AUTO_INCREMENT"),("reservation_id","INT FK → reservations.id"),
     ("mode","ENUM(carte,paypal,virement,crypto)"),
     ("montant","DECIMAL(10,2)"),
     ("statut","ENUM(en_attente,valide,refuse,rembourse)"),
     ("date_paiement","DATETIME"),
     ("reference_externe","VARCHAR(100)")],
    Inches(10.55), Inches(1.7), Inches(2.65), RGBColor(0xFF,0x69,0x69))

sr_table(s, "itineraires",
    [("id","INT PK AUTO_INCREMENT"),("utilisateur_id","INT FK → utilisateurs.id"),
     ("nom","VARCHAR(200)"),("destination_id","INT FK → destinations.id"),
     ("date_debut","DATE"),("date_fin","DATE"),
     ("nb_voyageurs","INT DEFAULT 1"),
     ("statut","ENUM(brouillon,confirme,annule)"),
     ("cout_total","DECIMAL(10,2)"),
     ("date_creation","DATETIME DEFAULT NOW()")],
    Inches(0.3), Inches(4.8), Inches(3.3), RGBColor(0x9B,0x59,0xB6))

sr_table(s, "itineraire_items",
    [("id","INT PK AUTO_INCREMENT"),("itineraire_id","INT FK → itineraires.id"),
     ("type_element","ENUM(transport,hebergement,activite,autre)"),
     ("element_ref_id","INT"),("quantite","INT DEFAULT 1"),
     ("prix_unitaire","DECIMAL(10,2)"),("date_debut","DATETIME"),
     ("date_fin","DATETIME"),("ordre","INT DEFAULT 0")],
    Inches(3.85), Inches(4.8), Inches(3.4), RGBColor(0x2E,0xCC,0x71))

page_num(s, 14)

# Slide 15 : SR Compléments
s = new_slide()
fill_slide_bg(s)
section_tag(s, "04 — SCHÉMA RELATIONNEL")
slide_title(s, "Schéma Relationnel — Compléments & Sécurité",
            "Tables : notifications · sessions · paniers · elements_panier · favoris · avis")
h_line(s, Inches(1.52))

sr_table(s, "notifications",
    [("id","INT PK AUTO_INCREMENT"),("utilisateur_id","INT FK → utilisateurs.id"),
     ("type","ENUM(confirmation,rappel,promo,alerte)"),
     ("titre","VARCHAR(200)"),("message","TEXT"),
     ("lu","BOOLEAN DEFAULT 0"),
     ("created_at","DATETIME DEFAULT NOW()")],
    Inches(0.3), Inches(1.7), Inches(3.2), ACCENT1)

sr_table(s, "sessions",
    [("id","INT PK AUTO_INCREMENT"),("utilisateur_id","INT FK → utilisateurs.id"),
     ("token","VARCHAR(255) UNIQUE NOT NULL"),
     ("expire_le","DATETIME NOT NULL"),
     ("ip_adresse","VARCHAR(45)"),
     ("created_at","DATETIME DEFAULT NOW()")],
    Inches(3.7), Inches(1.7), Inches(3.2), ACCENT2)

sr_table(s, "paniers",
    [("id","INT PK AUTO_INCREMENT"),("utilisateur_id","INT FK → utilisateurs.id"),
     ("date_mise_a_jour","DATETIME DEFAULT NOW()")],
    Inches(7.1), Inches(1.7), Inches(2.7), ACCENT3)

sr_table(s, "elements_panier",
    [("id","INT PK AUTO_INCREMENT"),("panier_id","INT FK → paniers.id"),
     ("type_element","ENUM(transport,hebergement,activite)"),
     ("element_ref_id","INT NOT NULL"),("quantite","INT DEFAULT 1"),
     ("prix_unitaire","DECIMAL(10,2)"),
     ("date_ajout","DATETIME DEFAULT NOW()")],
    Inches(10.05), Inches(1.7), Inches(3.15), RGBColor(0xFF,0x69,0x69))

sr_table(s, "favoris",
    [("id","INT PK AUTO_INCREMENT"),("utilisateur_id","INT FK → utilisateurs.id"),
     ("type_element","ENUM(destination,transport,hebergement,activite)"),
     ("element_ref_id","INT NOT NULL"),
     ("date_ajout","DATETIME DEFAULT NOW()")],
    Inches(0.3), Inches(4.8), Inches(3.2), RGBColor(0x9B,0x59,0xB6))

sr_table(s, "avis",
    [("id","INT PK AUTO_INCREMENT"),("utilisateur_id","INT FK → utilisateurs.id"),
     ("type_cible","ENUM(destination,transport,hebergement,activite)"),
     ("id_cible","INT NOT NULL"),("note","INT CHECK(note BETWEEN 1 AND 5)"),
     ("commentaire","TEXT"),
     ("date_avis","DATETIME DEFAULT NOW()")],
    Inches(3.7), Inches(4.8), Inches(3.2), RGBColor(0x2E,0xCC,0x71))

sr_table(s, "voyageurs_sejour",
    [("id","INT PK AUTO_INCREMENT"),("reservation_id","INT FK → reservations.id"),
     ("utilisateur_id","INT FK → utilisateurs.id"),
     ("refe_sejour","VARCHAR(50)"),
     ("date_ajout","DATETIME DEFAULT NOW()")],
    Inches(7.1), Inches(4.8), Inches(3.0), RGBColor(0xF3,0x9C,0x12))

page_num(s, 15)

# ─── SLIDE 16 : ARCHITECTURE — Vue globale ───────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "05 — ARCHITECTURE SYSTÈME")
slide_title(s, "Architecture Générale — Vue d'ensemble",
            "Frontend React · API PHP REST · MySQL · MAMP")
h_line(s, Inches(1.52))

# Trois colonnes : Frontend | Backend | BDD
cols = [
    (Inches(0.35), "FRONTEND", "React 18 + Vite", ACCENT1,
     ["React Router DOM v6", "Axios (HTTP Client)", "i18n FR/EN", "CSS Variables (dark/light)", "Vite Dev Server :3000"]),
    (Inches(4.65), "BACKEND", "PHP 8 · API REST", ACCENT2,
     ["Architecture MVC", "PDO + Prepared Statements", "Bcrypt Auth", "CORS Middleware", "Apache / MAMP :80"]),
    (Inches(8.95), "BASE DE DONNÉES", "MySQL 8", ACCENT3,
     ["14 tables relationnelles", "Clés étrangères (FK)", "Migrations versionnées", "utf8mb4 unicode", "Transactions ACID"]),
]
for x, title, sub, color, items in cols:
    w = Inches(4.0)
    add_rect(s, x, Inches(1.7), w, Inches(4.8), CARD_BG)
    add_rect(s, x, Inches(1.7), w, Inches(0.06), color)
    add_textbox(s, title, x + Inches(0.15), Inches(1.78),
                w - Inches(0.3), Inches(0.4), size=14, bold=True, color=color)
    add_textbox(s, sub, x + Inches(0.15), Inches(2.18),
                w - Inches(0.3), Inches(0.32), size=10, color=LIGHT_GREY, italic=True)
    add_rect(s, x + Inches(0.15), Inches(2.5), w - Inches(0.3), Inches(0.02), MID_GREY)
    for j, item in enumerate(items):
        add_textbox(s, "▸  " + item,
                    x + Inches(0.2), Inches(2.6) + j * Inches(0.5),
                    w - Inches(0.3), Inches(0.44), size=11, color=LIGHT_GREY)

# Arrows between cols
for x1, x2 in [(4.35, 4.65), (8.65, 8.95)]:
    arrow(s, x1, 4.1, x2, 4.1, "", ACCENT1)
    add_textbox(s, "HTTP\n/api", Inches((x1+x2)/2 - 0.18), Inches(3.82),
                Inches(0.5), Inches(0.4), size=8, color=ACCENT1, align=PP_ALIGN.CENTER)

# Communication note
add_rect(s, Inches(0.35), Inches(6.6), Inches(12.7), Inches(0.65), CARD_BG)
txb = s.shapes.add_textbox(Inches(0.5), Inches(6.62), Inches(12.4), Inches(0.6))
tf = txb.text_frame; tf.word_wrap = True
add_para(tf, "Flux : Browser → Vite proxy /api → Apache localhost:80 → PHP Controller → Model → PDO → MySQL",
         size=10, color=LIGHT_GREY)

page_num(s, 16)

# ─── SLIDE 17 : ARCHITECTURE — Flux données ──────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "05 — ARCHITECTURE SYSTÈME")
slide_title(s, "Architecture — Flux de Données & Interactions",
            "Détail des couches : Routing · Services · Controllers · Models")
h_line(s, Inches(1.52))

layers = [
    (Inches(0.3),  "COUCHE PRÉSENTATION",  ACCENT1,
     "React 18 · Vite · React Router v6",
     ["9 Screens (ScreenHome, ScreenResults, ScreenDetail,\n   ScreenCart, ScreenPayment, ScreenAccount,\n   ScreenItinerary, ScreenTransport, ScreenAdmin)",
      "8 Composants réutilisables (Header, Footer,\n   AuthModal, SearchBar, Notifications, Toast…)",
      "Services Axios : authService, destinationsService,\n   reservationsService, itinerairesService…"]),
    (Inches(4.6),  "COUCHE API / CONTRÔLEURS",  ACCENT2,
     "PHP 8 · Architecture MVC",
     ["9 endpoints REST (/auth, /destinations, /hebergements,\n   /activites, /transports, /reservations,\n   /itineraires, /notifications, /admin)",
      "Controllers : AuthController, DestinationController,\n   ReservationController, ItineraireController…",
      "Middleware : cors.php (CORS headers),\n   Validation entrées, Gestion erreurs HTTP"]),
    (Inches(8.9),  "COUCHE DONNÉES",  ACCENT3,
     "PHP PDO · MySQL 8",
     ["7 Models : Destination, Hebergement, Activite,\n   Transport, Reservation, Itineraire, Notification",
      "Transactions : décrémentation stocks (places,\n   chambres) avec rollback en cas d'annulation",
      "Requêtes préparées PDO, gestion\n   des relations N:M (reservation_activites)"]),
]
for x, title, color, sub, items in layers:
    w = Inches(4.1)
    add_rect(s, x, Inches(1.7), w, Inches(5.35), CARD_BG)
    add_rect(s, x, Inches(1.7), w, Inches(0.06), color)
    add_textbox(s, title, x + Inches(0.15), Inches(1.78),
                w - Inches(0.3), Inches(0.38), size=12, bold=True, color=color)
    add_textbox(s, sub, x + Inches(0.15), Inches(2.16),
                w - Inches(0.3), Inches(0.3), size=9.5, color=LIGHT_GREY, italic=True)
    add_rect(s, x + Inches(0.15), Inches(2.47), w - Inches(0.3), Inches(0.02), MID_GREY)
    for j, item in enumerate(items):
        add_textbox(s, "▸  " + item,
                    x + Inches(0.2), Inches(2.56) + j * Inches(1.45),
                    w - Inches(0.3), Inches(1.3), size=10, color=LIGHT_GREY)

for x1, x2 in [(4.4, 4.6), (8.7, 8.9)]:
    arrow(s, x1, 4.4, x2, 4.4, "", ACCENT2)

page_num(s, 17)

# ─── SLIDE 18 : ARCHITECTURE — Flux réservation ──────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "05 — ARCHITECTURE SYSTÈME")
slide_title(s, "Architecture — Flux de Réservation Complet",
            "Séquence de bout en bout : Recherche → Paiement → Confirmation")
h_line(s, Inches(1.52))

steps = [
    ("1", "Recherche", "L'utilisateur entre destination,\ndates et nb voyageurs dans SearchBar", ACCENT1),
    ("2", "Résultats", "ScreenResults filtre et affiche\nles destinations correspondantes", ACCENT2),
    ("3", "Détail", "ScreenDetail : onglets Transport,\nHébergement, Activités — sélection", ACCENT3),
    ("4", "Panier", "Ajout au panier (localStorage),\nrécapitulatif dans ScreenCart", RGBColor(0xFF,0x69,0x69)),
    ("5", "Réservation", "POST /reservations : validation dates,\nvérification stocks, génération ref VV-XXXXX", RGBColor(0x9B,0x59,0xB6)),
    ("6", "Paiement", "Création paiement, décrémentation\nstocks (places, chambres)", RGBColor(0x2E,0xCC,0x71)),
    ("7", "Notification", "Création notification confirmation,\nredirection ScreenAccount", RGBColor(0xF3,0x9C,0x12)),
]
for i, (num, title, desc, color) in enumerate(steps):
    col = i % 4
    row = i // 4
    x = Inches(0.35) + col * Inches(3.22)
    y = Inches(1.72) + row * Inches(2.4)
    w, h = Inches(3.0), Inches(2.15)
    add_rect(s, x, y, w, h, CARD_BG)
    add_rect(s, x, y, Inches(0.55), h, color)
    add_textbox(s, num, x + Inches(0.08), y + Inches(0.8),
                Inches(0.38), Inches(0.52), size=20, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, title, x + Inches(0.65), y + Inches(0.12),
                Inches(2.25), Inches(0.4), size=12, bold=True, color=color)
    add_textbox(s, desc, x + Inches(0.65), y + Inches(0.55),
                Inches(2.25), Inches(1.4), size=9.5, color=LIGHT_GREY)
    # Arrow to next
    if i < len(steps) - 1 and col < 3:
        arrow(s, (x + w)/Inches(1), (y + h/2)/Inches(1),
              (x + w + Inches(0.22))/Inches(1), (y + h/2)/Inches(1),
              "", ACCENT1)

page_num(s, 18)

# ─── SLIDE 19 : ARCHITECTURE — Rôles & Sécurité ──────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "05 — ARCHITECTURE SYSTÈME")
slide_title(s, "Rôles Utilisateurs & Sécurité",
            "3 rôles distincts · Authentification bcrypt · Gestion des droits")
h_line(s, Inches(1.52))

roles = [
    ("VISITEUR", "Non authentifié", ACCENT1,
     ["Consultation des destinations", "Recherche et filtrage", "Consultation des offres", "Inscription / Connexion"]),
    ("UTILISATEUR", "Authentifié (role: user)", ACCENT2,
     ["Toutes actions Visiteur", "Réservation & paiement", "Gestion espace compte", "Création d'itinéraires", "Favoris & avis", "Historique réservations", "Notifications personnalisées"]),
    ("ADMINISTRATEUR", "Authentifié (role: admin)", ACCENT3,
     ["Toutes actions Utilisateur", "Dashboard admin", "Gestion utilisateurs (CRUD)", "Modification rôles", "Accès données globales"]),
]
for i, (role, sub, color, perms) in enumerate(roles):
    x = Inches(0.35) + i * Inches(4.28)
    w = Inches(4.0)
    add_rect(s, x, Inches(1.72), w, Inches(4.9), CARD_BG)
    add_rect(s, x, Inches(1.72), w, Inches(0.06), color)
    add_textbox(s, role, x + Inches(0.2), Inches(1.8),
                w - Inches(0.3), Inches(0.4), size=14, bold=True, color=color)
    add_textbox(s, sub, x + Inches(0.2), Inches(2.2),
                w - Inches(0.3), Inches(0.3), size=10, color=LIGHT_GREY, italic=True)
    add_rect(s, x + Inches(0.2), Inches(2.52), w - Inches(0.4), Inches(0.02), MID_GREY)
    for j, perm in enumerate(perms):
        add_textbox(s, "✓  " + perm,
                    x + Inches(0.2), Inches(2.62) + j * Inches(0.46),
                    w - Inches(0.3), Inches(0.4), size=11, color=LIGHT_GREY)

# Sécurité note
add_rect(s, Inches(0.35), Inches(6.75), Inches(12.7), Inches(0.55), CARD_BG)
add_textbox(s,
    "Sécurité : Mots de passe hashés bcrypt · Rôle vérifié côté serveur (header X-User-Role) · PDO + requêtes préparées (anti-SQLi) · CORS restreint à localhost:3000",
    Inches(0.5), Inches(6.78), Inches(12.4), Inches(0.48),
    size=9.5, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

page_num(s, 19)

# ─── SLIDE 20 : SPECS FONC — Rôles & Périmètre ───────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "06 — SPÉCIFICATIONS FONCTIONNELLES")
slide_title(s, "Spécifications — Périmètre & Rôles",
            "Présentation générale de VoyageVista et de ses acteurs")
h_line(s, Inches(1.52))

# Présentation projet
add_rect(s, Inches(0.35), Inches(1.72), Inches(12.7), Inches(1.5), CARD_BG)
add_rect(s, Inches(0.35), Inches(1.72), Inches(0.08), Inches(1.5), ACCENT1)
add_textbox(s, "VoyageVista est une plateforme de voyage en ligne permettant aux utilisateurs de "
               "rechercher des destinations, composer des séjours sur-mesure (transport + hébergement + activités), "
               "créer des itinéraires personnalisés et effectuer des réservations en ligne. "
               "L'interface est disponible en français et en anglais avec support du mode sombre.",
            Inches(0.6), Inches(1.82), Inches(12.2), Inches(1.25),
            size=11, color=LIGHT_GREY)

roles2 = [
    ("Visiteur", ACCENT1, "Accès public : consultation, recherche, inscription"),
    ("Utilisateur", ACCENT2, "Réservation, itinéraires, compte, favoris, avis"),
    ("Administrateur", ACCENT3, "Dashboard, gestion utilisateurs, données"),
]
for i, (role, color, desc) in enumerate(roles2):
    x = Inches(0.35) + i * Inches(4.28)
    w = Inches(4.0)
    add_rect(s, x, Inches(3.4), w, Inches(1.0), CARD_BG)
    add_rect(s, x, Inches(3.4), Inches(0.08), Inches(1.0), color)
    add_textbox(s, role, x + Inches(0.2), Inches(3.48),
                w - Inches(0.25), Inches(0.38), size=13, bold=True, color=color)
    add_textbox(s, desc, x + Inches(0.2), Inches(3.86),
                w - Inches(0.25), Inches(0.45), size=10, color=LIGHT_GREY)

# Contraintes techniques
add_rect(s, Inches(0.35), Inches(4.6), Inches(12.7), Inches(2.6), CARD_BG)
add_rect(s, Inches(0.35), Inches(4.6), Inches(0.08), Inches(2.6), ACCENT3)
add_textbox(s, "Contraintes Techniques",
            Inches(0.6), Inches(4.67), Inches(6), Inches(0.38),
            size=13, bold=True, color=ACCENT3)
constraints = [
    ("Performance", "Chargement < 3s, pagination des listes, lazy loading images"),
    ("Sécurité", "HTTPS, anti-injection SQL (PDO), hash bcrypt, CORS restreint"),
    ("Compatibilité", "Navigateurs modernes (Chrome, Firefox, Safari, Edge), responsive mobile"),
    ("Disponibilité", "Architecture stateless, séparation frontend/backend déployable indépendamment"),
    ("Internationalisation", "FR / EN intégré, extensible à d'autres langues via i18n.js"),
    ("Accessibilité", "Contraste WCAG AA, navigation clavier, ARIA labels"),
]
for i, (key, val) in enumerate(constraints):
    col = i % 2; row = i // 2
    x = Inches(0.6) + col * Inches(6.3)
    y = Inches(5.1) + row * Inches(0.55)
    add_textbox(s, key + " : ", x, y, Inches(1.5), Inches(0.45),
                size=10, bold=True, color=ACCENT1)
    add_textbox(s, val, x + Inches(1.5), y, Inches(4.65), Inches(0.45),
                size=10, color=LIGHT_GREY)

page_num(s, 20)

# ─── SLIDE 21 : SPECS — Catalogue & Recherche ────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "06 — SPÉCIFICATIONS FONCTIONNELLES")
slide_title(s, "SF-01 · Catalogue & Recherche",
            "Consultation des destinations et fonctionnalités de recherche")
h_line(s, Inches(1.52))

feats = [
    ("SF-01.1", "Affichage Catalogue", ACCENT1,
     "Afficher la liste paginée des destinations actives avec photo, nom, pays, note, "
     "prix depuis, tag catégorie. Filtrage par catégorie (Asie, Europe, Amériques…)."),
    ("SF-01.2", "Recherche Globale", ACCENT2,
     "Barre de recherche combinée : destination + dates aller/retour + nombre de voyageurs. "
     "Résultats filtrés en temps réel sur ScreenResults."),
    ("SF-01.3", "Filtres Avancés", ACCENT3,
     "Sur la page résultats : filtres par budget, note minimale, type de destination, "
     "durée de séjour. Tri par pertinence / prix / note."),
    ("SF-01.4", "Détail Destination", RGBColor(0xFF,0x69,0x69),
     "Page détail avec 4 onglets : Vue d'ensemble (description, galerie), "
     "Transports disponibles, Hébergements, Activités. Pré-sélection en un clic."),
    ("SF-01.5", "Recherche Transport", RGBColor(0x9B,0x59,0xB6),
     "Page dédiée ScreenTransport : recherche multimodale (avion, train, bus, voiture) "
     "avec filtres type, compagnie, prix min/max, places disponibles."),
    ("SF-01.6", "Internationalisation", RGBColor(0x2E,0xCC,0x71),
     "Toutes les destinations disposent de titres, descriptions et tags en FR et EN. "
     "Toggle langue en header, persisté en localStorage."),
]
for i, (ref, title, color, desc) in enumerate(feats):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.48)
    y = Inches(1.72) + row * Inches(1.72)
    w = Inches(6.2)
    add_rect(s, x, y, w, Inches(1.55), CARD_BG)
    add_rect(s, x, y, Inches(0.08), Inches(1.55), color)
    add_textbox(s, ref, x + Inches(0.18), y + Inches(0.06),
                Inches(1.0), Inches(0.3), size=9, color=color, italic=True)
    add_textbox(s, title, x + Inches(1.2), y + Inches(0.06),
                Inches(4.9), Inches(0.3), size=11, bold=True, color=WHITE)
    add_textbox(s, desc, x + Inches(0.18), y + Inches(0.4),
                w - Inches(0.3), Inches(1.05), size=9.5, color=LIGHT_GREY)

page_num(s, 21)

# ─── SLIDE 22 : SPECS — Réservation & Panier ─────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "06 — SPÉCIFICATIONS FONCTIONNELLES")
slide_title(s, "SF-02 · Réservation & Panier",
            "Composition de séjour, gestion du panier et paiement")
h_line(s, Inches(1.52))

feats2 = [
    ("SF-02.1", "Composition de Séjour", ACCENT1,
     "Depuis ScreenDetail, l'utilisateur sélectionne transport, hébergement et activités "
     "souhaitées. Chaque élément est ajouté au panier (stocké en localStorage)."),
    ("SF-02.2", "Gestion du Panier", ACCENT2,
     "ScreenCart : liste des éléments sélectionnés, modification quantité, suppression "
     "d'éléments, calcul du total en temps réel. Badge compteur dans le header."),
    ("SF-02.3", "Validation Réservation", ACCENT3,
     "POST /reservations : vérification disponibilités (stocks), validation des dates "
     "(date_retour > date_depart), génération référence unique VV-XXXXXXX."),
    ("SF-02.4", "Gestion Activités", RGBColor(0xFF,0x69,0x69),
     "Sélection multiple d'activités liées à la destination. Enregistrement dans "
     "reservation_activites (table N:M), décrémentation places_restantes."),
    ("SF-02.5", "Paiement", RGBColor(0x9B,0x59,0xB6),
     "ScreenPayment : choix du mode de paiement (carte, PayPal, virement, crypto), "
     "récapitulatif final, confirmation de commande."),
    ("SF-02.6", "Annulation & Rollback", RGBColor(0x2E,0xCC,0x71),
     "Annulation depuis espace compte : remise à jour automatique des stocks "
     "(nb_chambres_dispo, places_dispo, places_restantes) avec transaction ACID."),
]
for i, (ref, title, color, desc) in enumerate(feats2):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.48)
    y = Inches(1.72) + row * Inches(1.72)
    w = Inches(6.2)
    add_rect(s, x, y, w, Inches(1.55), CARD_BG)
    add_rect(s, x, y, Inches(0.08), Inches(1.55), color)
    add_textbox(s, ref, x + Inches(0.18), y + Inches(0.06),
                Inches(1.0), Inches(0.3), size=9, color=color, italic=True)
    add_textbox(s, title, x + Inches(1.2), y + Inches(0.06),
                Inches(4.9), Inches(0.3), size=11, bold=True, color=WHITE)
    add_textbox(s, desc, x + Inches(0.18), y + Inches(0.4),
                w - Inches(0.3), Inches(1.05), size=9.5, color=LIGHT_GREY)

page_num(s, 22)

# ─── SLIDE 23 : SPECS — Itinéraires & Compte ─────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "06 — SPÉCIFICATIONS FONCTIONNELLES")
slide_title(s, "SF-03 · Itinéraires & Espace Compte",
            "Planification personnalisée et gestion du profil utilisateur")
h_line(s, Inches(1.52))

feats3 = [
    ("SF-03.1", "Créateur d'Itinéraire", ACCENT1,
     "ScreenItinerary : créer un itinéraire nommé avec destination, dates. "
     "Ajouter des éléments ordonnés (transport, hébergement, activité, autre) "
     "avec quantité, prix unitaire, plage horaire."),
    ("SF-03.2", "Sauvegarde & Statuts", ACCENT2,
     "Statuts : brouillon (sauvegarde partielle), confirmé, annulé. "
     "Calcul automatique du coût total. API PUT /itineraires pour mise à jour."),
    ("SF-03.3", "Gestion du Profil", ACCENT3,
     "ScreenAccount : modification nom, prénom, email, téléphone, date de naissance. "
     "Changement de mot de passe avec vérification de l'ancien."),
    ("SF-03.4", "Historique Réservations", RGBColor(0xFF,0x69,0x69),
     "Liste des réservations avec statut (en attente, confirmée, annulée, terminée), "
     "montant, référence, dates. Détail de chaque réservation accessible."),
    ("SF-03.5", "Notifications", RGBColor(0x9B,0x59,0xB6),
     "Notifications temps réel pour : confirmation réservation, rappels, "
     "promotions, alertes. Badge non-lu dans le header, marquage lu individuel/global."),
    ("SF-03.6", "Favoris & Avis", RGBColor(0x2E,0xCC,0x71),
     "Ajout en favoris de destinations, transports, hébergements, activités. "
     "Dépôt d'avis notés (1-5 étoiles) avec commentaire sur tout élément."),
]
for i, (ref, title, color, desc) in enumerate(feats3):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.48)
    y = Inches(1.72) + row * Inches(1.72)
    w = Inches(6.2)
    add_rect(s, x, y, w, Inches(1.55), CARD_BG)
    add_rect(s, x, y, Inches(0.08), Inches(1.55), color)
    add_textbox(s, ref, x + Inches(0.18), y + Inches(0.06),
                Inches(1.0), Inches(0.3), size=9, color=color, italic=True)
    add_textbox(s, title, x + Inches(1.2), y + Inches(0.06),
                Inches(4.9), Inches(0.3), size=11, bold=True, color=WHITE)
    add_textbox(s, desc, x + Inches(0.18), y + Inches(0.4),
                w - Inches(0.3), Inches(1.05), size=9.5, color=LIGHT_GREY)

page_num(s, 23)

# ─── SLIDE 24 : SPECS — Administration ───────────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "06 — SPÉCIFICATIONS FONCTIONNELLES")
slide_title(s, "SF-04 · Administration",
            "Dashboard admin — Gestion des utilisateurs et des données")
h_line(s, Inches(1.52))

feats4 = [
    ("SF-04.1", "Dashboard Administrateur", ACCENT1,
     "ScreenAdmin accessible uniquement si role === 'admin'. Vue d'ensemble des "
     "utilisateurs, réservations, statistiques générales de la plateforme."),
    ("SF-04.2", "Gestion Utilisateurs", ACCENT2,
     "Liste complète des utilisateurs avec nom, email, rôle, statut actif. "
     "Changement de rôle (user ↔ admin), suppression de compte (DELETE /admin)."),
    ("SF-04.3", "Gestion Destinations", ACCENT3,
     "Création (POST /destinations), modification (PUT) et suppression (DELETE) "
     "de destinations. Champ actif pour masquer sans supprimer."),
    ("SF-04.4", "Gestion Offres", RGBColor(0xFF,0x69,0x69),
     "CRUD complet sur transports, hébergements, activités via endpoints dédiés. "
     "Contrôle des stocks et disponibilités depuis l'interface admin."),
]
for i, (ref, title, color, desc) in enumerate(feats4):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.48)
    y = Inches(1.72) + row * Inches(1.9)
    w = Inches(6.2)
    add_rect(s, x, y, w, Inches(1.72), CARD_BG)
    add_rect(s, x, y, Inches(0.08), Inches(1.72), color)
    add_textbox(s, ref, x + Inches(0.18), y + Inches(0.07),
                Inches(1.0), Inches(0.32), size=9, color=color, italic=True)
    add_textbox(s, title, x + Inches(1.2), y + Inches(0.07),
                Inches(4.9), Inches(0.32), size=11, bold=True, color=WHITE)
    add_textbox(s, desc, x + Inches(0.18), y + Inches(0.45),
                w - Inches(0.3), Inches(1.15), size=10, color=LIGHT_GREY)

# Techno stack
add_rect(s, Inches(0.35), Inches(5.75), Inches(12.7), Inches(1.5), CARD_BG)
add_rect(s, Inches(0.35), Inches(5.75), Inches(0.08), Inches(1.5), RGBColor(0x9B,0x59,0xB6))
add_textbox(s, "Technologies Envisagées",
            Inches(0.6), Inches(5.82), Inches(12.2), Inches(0.38),
            size=13, bold=True, color=RGBColor(0x9B,0x59,0xB6))
techs = [
    ("Frontend", "React 18, Vite 5, React Router v6, Axios, CSS Variables"),
    ("Backend", "PHP 8, Architecture MVC, PDO, Bcrypt, Apache/MAMP"),
    ("Base de données", "MySQL 8, utf8mb4, Transactions ACID, Migrations versionnées"),
    ("Dev & Outillage", "Git, MAMP (local), Vite proxy, ESLint, Google Fonts"),
]
for i, (key, val) in enumerate(techs):
    col = i % 2; row = i // 2
    x = Inches(0.6) + col * Inches(6.35)
    y = Inches(6.28) + row * Inches(0.45)
    add_textbox(s, key + " : ", x, y, Inches(1.7), Inches(0.38),
                size=10, bold=True, color=ACCENT1)
    add_textbox(s, val, x + Inches(1.7), y, Inches(4.5), Inches(0.38),
                size=10, color=LIGHT_GREY)

page_num(s, 24)

# ─── SLIDE 25 : SPECS — Cohérence & Synthèse ─────────────────────────────────
s = new_slide()
fill_slide_bg(s)
section_tag(s, "06 — SPÉCIFICATIONS FONCTIONNELLES")
slide_title(s, "Synthèse des Fonctionnalités",
            "Vue consolidée — correspondance fonctionnalités / modules techniques")
h_line(s, Inches(1.52))

rows = [
    ("SF-01", "Catalogue & Recherche", "ScreenHome, ScreenResults, ScreenDetail, ScreenTransport", "destinations, transports, hebergements, activites", ACCENT1),
    ("SF-02", "Réservation & Panier", "ScreenCart, ScreenPayment", "reservations, reservation_activites, paiements", ACCENT2),
    ("SF-03.1", "Itinéraires", "ScreenItinerary", "itineraires, itineraire_items", ACCENT3),
    ("SF-03.3", "Espace Compte", "ScreenAccount", "utilisateurs, sessions, notifications", RGBColor(0xFF,0x69,0x69)),
    ("SF-03.6", "Favoris & Avis", "ScreenAccount (onglets)", "favoris, avis", RGBColor(0x9B,0x59,0xB6)),
    ("SF-04", "Administration", "ScreenAdmin", "utilisateurs (admin CRUD)", RGBColor(0x2E,0xCC,0x71)),
]

# Header
hdrs = [("Réf.", 0.6), ("Fonctionnalité", 1.8), ("Écran(s) React", 3.5), ("Table(s) MySQL", 3.8)]
x = Inches(0.35)
y = Inches(1.72)
for label, w in hdrs:
    add_rect(s, x, y, Inches(w), Inches(0.4), ACCENT1)
    add_textbox(s, label, x + Inches(0.08), y + Inches(0.04),
                Inches(w - 0.15), Inches(0.32), size=10, bold=True, color=TEXT_DARK)
    x += Inches(w)

for i, (ref, feat, screens, tables, color) in enumerate(rows):
    y = Inches(2.18) + i * Inches(0.72)
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x12,0x20,0x30)
    x = Inches(0.35)
    for j, (val, w) in enumerate([(ref, 0.6), (feat, 1.8), (screens, 3.5), (tables, 3.8)]):
        add_rect(s, x, y, Inches(w), Inches(0.65), bg)
        c = color if j == 0 else (WHITE if j == 1 else LIGHT_GREY)
        add_textbox(s, val, x + Inches(0.08), y + Inches(0.06),
                    Inches(w - 0.15), Inches(0.55), size=9 if j > 1 else 10,
                    bold=(j == 1), color=c)
        x += Inches(w)

page_num(s, 25)

# ─── SAVE ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\wanga\Downloads\VoyageVista_PowerPoint_Preliminaire.pptx"
prs.save(out)
print("SAVED:", out)
print("Total slides:", len(prs.slides))
